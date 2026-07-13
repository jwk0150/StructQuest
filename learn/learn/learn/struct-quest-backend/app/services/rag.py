import os
import uuid
import time
import json
from typing import List, Optional, Dict

import httpx

from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

from pptx import Presentation as PptxPresentation

# ChromaDB 持久化路径
CHROMA_PATH = os.getenv("CHROMA_DB_PATH", "./chroma_db")

# Embedding API 配置（SiliconFlow 硅基流动，国内直连，免费）
EMBEDDING_API_KEY = os.getenv("EMBEDDING_API_KEY", "")
EMBEDDING_BASE_URL = os.getenv("EMBEDDING_BASE_URL", "https://api.siliconflow.cn/v1")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "BAAI/bge-large-zh-v1.5")

# 切分配置
CHUNK_SIZE = int(os.getenv("RAG_CHUNK_SIZE", "1000"))
CHUNK_OVERLAP = int(os.getenv("RAG_CHUNK_OVERLAP", "200"))


class EmbeddingClient:
    """直接调用 SiliconFlow Embedding API（无需 langchain-openai）"""

    # 单次请求最大 token 数（bge-m3 支持 8192，但保险起见用 4000）
    MAX_TOKENS_PER_REQUEST = 4000
    # 粗略估算：1 个中文字 ≈ 1.5 tokens，1 个英文单词 ≈ 1.3 tokens
    ESTIMATED_CHARS_PER_TOKEN = 1.5

    def __init__(self):
        if not EMBEDDING_API_KEY:
            print("[EmbeddingClient] 警告: EMBEDDING_API_KEY 未配置！请检查 .env 文件")
        self._client = httpx.Client(
            base_url=EMBEDDING_BASE_URL,
            headers={
                "Authorization": f"Bearer {EMBEDDING_API_KEY}",
                "Content-Type": "application/json",
            },
            timeout=60.0,
        )

    @staticmethod
    def _estimate_tokens(text: str) -> int:
        """粗略估算文本的 token 数"""
        # 中文字符数 * 1.5 + 非中文字符数 / 4（英文平均每词 4 字符）
        cn_chars = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        other_chars = len(text) - cn_chars
        return int(cn_chars * 1.5 + other_chars / 4)

    @staticmethod
    def _split_long_text(text: str, max_chars: int) -> List[str]:
        """将超长文本按句子边界拆分"""
        if len(text) <= max_chars:
            return [text]
        parts = []
        remaining = text
        while remaining:
            if len(remaining) <= max_chars:
                parts.append(remaining)
                break
            # 在 max_chars 附近找句号、换行或空格作为分割点
            cut_pos = max_chars
            for sep in ['。\n', '。', '!\n', '!', '?\n', '?', '\n', ' ', '；', '，']:
                pos = remaining.rfind(sep, max_chars - 200, max_chars + 50)
                if pos > 0:
                    cut_pos = pos + len(sep)
                    break
            parts.append(remaining[:cut_pos])
            remaining = remaining[cut_pos:]
        return parts

    def _embed_batch(self, texts: List[str]) -> List[Dict]:
        """发送一个批次到 API，返回原始结果列表"""
        resp = self._client.post(
            "/embeddings",
            json={"model": EMBEDDING_MODEL, "input": texts},
        )
        if resp.status_code != 200:
            print(f"[EmbeddingClient] 错误状态码: {resp.status_code}")
            print(f"[EmbeddingClient] 错误响应: {resp.text}")
        resp.raise_for_status()
        return resp.json()["data"]

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """批量向量化文档（自动处理超长文本和批量限制）"""
        all_embeddings: List[List[float]] = []
        all_indices: List[int] = []  # 记录每个 embedding 对应原 texts 的索引

        batch_texts: List[str] = []
        batch_token_count = 0

        for idx, text in enumerate(texts):
            est_tokens = self._estimate_tokens(text)

            if est_tokens > self.MAX_TOKENS_PER_REQUEST:
                # 单条文本超长 → 拆分后分别嵌入，再取平均
                sub_parts = self._split_long_text(text, int(self.MAX_TOKENS_PER_REQUEST / self.ESTIMATED_CHARS_PER_TOKEN))
                sub_embeddings = []
                for part in sub_parts:
                    result = self._embed_batch([part])
                    sub_embeddings.append(result[0]["embedding"])
                # 取所有子段的平均值作为整体嵌入
                avg_emb = []
                if sub_embeddings:
                    dim = len(sub_embeddings[0])
                    for d in range(dim):
                        avg_emb.append(sum(e[d] for e in sub_embeddings) / len(sub_embeddings))
                all_embeddings.append(avg_emb)
                all_indices.append(idx)

            elif batch_token_count + est_tokens > self.MAX_TOKENS_PER_REQUEST:
                # 当前批次满了 → 先发送
                results = self._embed_batch(batch_texts)
                for r in sorted(results, key=lambda x: x["index"]):
                    all_embeddings.append(r["embedding"])
                    all_indices.append(batch_texts.pop(0))  # 占位
                # 开始新批次
                batch_texts = [text]
                batch_token_count = est_tokens
            else:
                batch_texts.append(text)
                batch_token_count += est_tokens

        # 发送剩余批次
        if batch_texts:
            results = self._embed_batch(batch_texts)
            for r in sorted(results, key=lambda x: x["index"]):
                all_embeddings.append(r["embedding"])

        return all_embeddings

    def embed_query(self, text: str) -> List[float]:
        """单条查询向量化"""
        try:
            est_tokens = self._estimate_tokens(text)
            if est_tokens > self.MAX_TOKENS_PER_REQUEST:
                parts = self._split_long_text(text, int(self.MAX_TOKENS_PER_REQUEST / self.ESTIMATED_CHARS_PER_TOKEN))
                sub_embeddings = []
                for part in parts:
                    result = self._embed_batch([part])
                    sub_embeddings.append(result[0]["embedding"])
                dim = len(sub_embeddings[0]) if sub_embeddings else 0
                return [sum(e[d] for e in sub_embeddings) / len(sub_embeddings) for d in range(dim)] if sub_embeddings else []

            result = self._embed_batch([text])
            return result[0]["embedding"]
        except Exception as e:
            print(f"[EmbeddingClient] embed_query 失败: {e}")
            raise


class RAGService:
    """RAG 知识库服务：PDF 处理、向量化存储、语义检索"""

    def __init__(self):
        self._embeddings: Optional[EmbeddingClient] = None
        self._text_splitter: Optional[RecursiveCharacterTextSplitter] = None
        self._collection = None
        self._chroma_client = None

    # ==================== 嵌入模型 =====================

    def _get_embeddings(self) -> EmbeddingClient:
        """懒初始化 Embedding 客户端"""
        if self._embeddings is None:
            print("[RAG] 正在初始化 DeepSeek Embedding 客户端 ...")
            self._embeddings = EmbeddingClient()
            print("[RAG] Embedding 客户端就绪!")
        return self._embeddings

    # ==================== ChromaDB 向量库 =====================

    def _get_collection(self):
        """懒初始化 ChromaDB 集合"""
        if self._collection is None:
            print(f"[RAG] 初始化向量数据库 (ChromaDB) ... 路径: {os.path.abspath(CHROMA_PATH)}")
            t0 = time.time()
            import chromadb
            self._chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)
            self._collection = self._chroma_client.get_or_create_collection(
                name="knowledge_base",
                metadata={"hnsw:space": "cosine"},
            )
            print(f"[RAG] 向量数据库就绪! 文档数: {self._collection.count()}, 耗时: {time.time()-t0:.1f}秒")
        return self._collection

    # ==================== 文本切分器 =====================

    @property
    def text_splitter(self) -> RecursiveCharacterTextSplitter:
        if self._text_splitter is None:
            self._text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=CHUNK_SIZE,
                chunk_overlap=CHUNK_OVERLAP,
                separators=["\n\n", "\n", "。", "！", "？", ".", " ", ""],
            )
        return self._text_splitter

    # ==================== 文档处理 =====================

    # ==================== 公共：切分+向量化+入库 =====================

    def _ingest_documents(
        self, docs: List[Document], doc_id: str, label: str
    ) -> Dict:
        """将 Document 列表切分、向量化并写入 ChromaDB"""
        total_start = time.time()

        # 切分
        print(f"[RAG][{doc_id}] 切分文档 ({label}) ...")
        t0 = time.time()
        splits = self.text_splitter.split_documents(docs)
        for i, split in enumerate(splits):
            split.metadata["doc_id"] = doc_id
            split.metadata["chunk_index"] = i
        print(f"[RAG][{doc_id}] 切分完成! 共 {len(splits)} 个片段, 耗时: {time.time()-t0:.1f}秒")

        if not splits:
            raise ValueError("切分后没有有效内容")

        # 向量化 + 入库
        print(f"[RAG][{doc_id}] 向量化并写入向量库 ({len(splits)} 片段) ...")
        t0 = time.time()

        collection = self._get_collection()

        documents = [split.page_content for split in splits]
        ids = [f"{doc_id}_{i}" for i in range(len(splits))]
        metadatas = []
        for split in splits:
            m = dict(split.metadata)
            for k, v in list(m.items()):
                if isinstance(v, int):
                    m[k] = str(v)
            metadatas.append(m)

        emb_client = self._get_embeddings()
        embeddings = emb_client.embed_documents(documents)

        batch_size = 50
        for i in range(0, len(documents), batch_size):
            end = min(i + batch_size, len(documents))
            collection.add(
                ids=ids[i:end],
                documents=documents[i:end],
                embeddings=embeddings[i:end],
                metadatas=metadatas[i:end],
            )
            print(f"[RAG][{doc_id}]   已写入 {end}/{len(documents)}")

        print(f"[RAG][{doc_id}] 全部完成! 总耗时: {time.time()-total_start:.1f}秒")

        return {
            "doc_id": doc_id,
            "chunks": len(splits),
            "total_chars": sum(len(d.page_content) for d in splits),
        }

    # ==================== 文档加载 =====================

    def process_pdf(self, file_path: str, doc_id: Optional[str] = None) -> Dict:
        """
        加载 PDF -> 文本切分 -> 向量化存入 ChromaDB
        返回处理结果字典
        """
        doc_id = doc_id or str(uuid.uuid4())[:8]

        print(f"[RAG][{doc_id}] Step1/2: 加载 PDF ...")
        t0 = time.time()
        try:
            loader = PyPDFLoader(file_path)
            docs = loader.load()
        except Exception as e:
            raise ValueError(f"PDF 加载失败: {e}") from e
        print(f"[RAG][{doc_id}] PDF 加载完成! 共 {len(docs)} 页, 耗时: {time.time()-t0:.1f}秒")

        if not docs:
            raise ValueError("PDF 文件为空或无法读取文本内容")

        return self._ingest_documents(docs, doc_id, "PDF")

    def process_pptx(self, file_path: str, doc_id: Optional[str] = None) -> Dict:
        """
        加载 PPTX -> 提取每页文本 -> 切分 -> 向量化存入 ChromaDB
        返回处理结果字典
        """
        doc_id = doc_id or str(uuid.uuid4())[:8]

        print(f"[RAG][{doc_id}] Step1/2: 加载 PPTX ...")
        t0 = time.time()
        try:
            prs = PptxPresentation(file_path)
            docs = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    page_content = "\n".join(texts)
                    docs.append(Document(
                        page_content=page_content,
                        metadata={"source": os.path.basename(file_path), "page": slide_num},
                    ))
        except Exception as e:
            raise ValueError(f"PPTX 加载失败: {e}") from e

        print(f"[RAG][{doc_id}] PPTX 加载完成! 共 {len(docs)} 张有效幻灯片, 耗时: {time.time()-t0:.1f}秒")

        if not docs:
            raise ValueError("PPTX 文件为空或无法提取文本内容")

        return self._ingest_documents(docs, doc_id, "PPTX")

    def delete_document(self, doc_id: str) -> int:
        """根据 doc_id 删除指定文档的所有向量块"""
        try:
            collection = self._get_collection()
            result = collection.get(where={"doc_id": doc_id})
            if not result or not result.get("ids"):
                return 0
            ids_to_delete = result["ids"]
            count = len(ids_to_delete)
            collection.delete(ids=ids_to_delete)
            print(f"[RAG] 删除文档 doc_id={doc_id}, 移除 {count} 个向量块")
            return count
        except Exception as e:
            print(f"[RAG] 删除失败: {e}")
            return -1

    # ==================== 语义检索 =====================

    def retrieve_context(self, query: str, k: int = 4) -> str:
        """语义检索：返回 top-k 相关上下文拼接后的文本"""
        if not query or not query.strip():
            return ""

        try:
            collection = self._get_collection()
            if collection.count() == 0:
                return ""

            emb_client = self._get_embeddings()
            query_embedding = emb_client.embed_query(query)

            results = collection.query(
                query_embeddings=[query_embedding],
                n_results=min(k, collection.count()),
                include=["documents", "metadatas"],
            )

            if not results or not results["documents"][0]:
                return ""

            context_parts = []
            docs = results["documents"][0]
            metas = results["metadatas"][0]
            for i, doc in enumerate(docs):
                meta = metas[i] if i < len(metas) else {}
                source = meta.get("source", "未知来源")
                page = meta.get("page", "?")
                context_parts.append(
                    f"[参考资料 {i + 1}](来源: {source}, 页码: {page})\n{doc}"
                )

            return "\n\n".join(context_parts)

        except Exception as e:
            print(f"[RAG] 检索错误: {e}")
            return ""

    def retrieve_with_scores(self, query: str, k: int = 4) -> List[Dict]:
        """语义检索：返回带相似度分数的结果列表"""
        results = []
        if not query or not query.strip():
            return results

        try:
            collection = self._get_collection()
            if collection.count() == 0:
                return results

            emb_client = self._get_embeddings()
            query_embedding = emb_client.embed_query(query)

            query_results = collection.query(
                query_embeddings=[query_embedding],
                n_results=min(k, collection.count()),
                include=["documents", "metadatas", "distances"],
            )

            if not query_results or not query_results["documents"][0]:
                return results

            docs = query_results["documents"][0]
            metas = query_results["metadatas"][0]
            dists = query_results["distances"][0]

            for i, doc in enumerate(docs):
                results.append({
                    "content": doc,
                    "metadata": metas[i] if i < len(metas) else {},
                    "score": round(float(dists[i]), 4) if i < len(dists) else 0,
                })
            return results
        except Exception as e:
            print(f"[RAG] 检索错误(带分数): {e}")
            return []

    # ==================== 统计与管理 ====================

    def get_stats(self) -> Dict:
        """获取知识库统计信息"""
        try:
            collection = self._get_collection()
            total_chunks = collection.count()

            all_meta = collection.get(include=["metadatas"])
            doc_ids = set()
            for m in (all_meta.get("metadatas") or []):
                if m and "doc_id" in m:
                    doc_ids.add(m["doc_id"])

            return {
                "total_chunks": total_chunks,
                "document_count": len(doc_ids),
                "chroma_path": CHROMA_PATH,
                "embedding_model": f"SiliconFlow ({EMBEDDING_MODEL})",
            }
        except Exception as e:
            return {"error": str(e)}

    def list_documents(self) -> List[Dict]:
        """列出所有已入库的文档及其统计"""
        try:
            collection = self._get_collection()
            all_meta = collection.get(include=["metadatas"])
            metadatas = all_meta.get("metadatas") or []

            doc_map: Dict[str, Dict] = {}
            for meta in metadatas:
                if not meta or "doc_id" not in meta:
                    continue
                did = meta["doc_id"]
                source = meta.get("source", "未知文件名")
                if did not in doc_map:
                    doc_map[did] = {
                        "doc_id": did,
                        "filename": os.path.basename(source) if source else "未知",
                        "source": source,
                        "chunks": 0,
                        "pages": set(),
                    }
                doc_map[did]["chunks"] += 1
                if "page" in meta:
                    doc_map[did]["pages"].add(meta["page"])

            result = list(doc_map.values())
            for item in result:
                item["page_count"] = len(item.pop("pages"))
                del item["pages"]

            return result
        except Exception as e:
            print(f"[RAG] 列表查询失败: {e}")
            return []


# 全局单例实例
rag_service = RAGService()
