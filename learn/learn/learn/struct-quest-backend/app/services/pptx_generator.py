"""
PPTX 演示文稿直接生成器（科大讯飞风格）
========================================

将 LLM 生成的 Markdown 大纲渲染为专业 .pptx 幻灯片。
设计灵感：现代企业级演示风格（彩色圆圈卡片、渐变、视觉层次）

依赖: python-pptx
"""

import os
import re
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List

# python-pptx 常用枚举（模块级导入，供默认参数使用）
try:
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Inches = Pt = Emu = None  # type: ignore
    RGBColor = None  # type: ignore
    PP_ALIGN = MSO_ANCHOR = MSO_SHAPE = None  # type: ignore

logger = logging.getLogger("pptx_generator")

def _get_output_dir() -> Path:
    static_base = os.environ.get("STATIC_DIR", "/app/static")
    pptx_dir = os.environ.get("PPTX_OUTPUT_DIR", os.path.join(static_base, "pptx"))
    dir_path = Path(pptx_dir)
    dir_path.mkdir(parents=True, exist_ok=True)
    return dir_path


# ── 科大讯飞风格主题配置 ──
THEME = {
    # 字体
    "title_font": "Microsoft YaHei",
    "body_font": "Microsoft YaHei",
    "code_font": "Consolas",
    "title_size": 44,
    "heading1_size": 32,
    "heading2_size": 24,
    "body_size": 16,
    "small_size": 13,
    "code_size": 13,

    # 品牌色系（参考科大讯飞）
    "brand_red": "E63946",        # 主品牌红
    "brand_blue": "1D3557",       # 深蓝
    "accent_orange": "F4A261",    # 强调橙
    "accent_teal": "2A9D8F",      # 青绿
    "accent_blue": "457B9D",      # 蓝色
    "accent_purple": "7209B7",    # 紫色
    "text_dark": "1D3557",        # 深色正文
    "text_gray": "455A64",        # 灰色文字
    "text_light": "78909C",       # 浅灰
    "bg_light": "F8F9FA",         # 浅背景
    "white": "FFFFFF",
    "card_shadow": "E0E0E0",

    # 圆圈配色循环（用于概念卡片）
    "circle_colors": [
        ("F4A261", "E76F51"),   # 橙 → 深橙
        ("E63946", "C1121F"),   # 红 → 暗红
        ("2A9D8F", "16817A"),   # 青 → 深青
        ("457B9D", "1D3557"),   # 蓝 → 深蓝
        ("7209B7", "560BAD"),   # 紫 → 深紫
        ("06D6A0", "059669"),   # 薄荷 → 绿
    ],
}


class PPTXGenerator:
    """PPTX 文件生成器 — 科大讯飞专业演示风格"""

    def __init__(self, output_dir: Optional[Path] = None):
        self.output_dir = output_dir or _get_output_dir()
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate(
        self,
        outline_content: str,
        topic: str = "学习内容",
        task_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        from datetime import datetime

        task_id = task_id or datetime.now().strftime("%Y%m%d%H%M%S%f")
        result = {
            "success": False,
            "file_url": None,
            "file_path": None,
            "slide_count": 0,
            "error": None,
            "download_name": f"{topic}.pptx",
        }

        try:
            from pptx import Presentation
            # Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE 已在模块级导入

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            slides_data = self._parse_markdown(outline_content, topic)

            for idx, slide_info in enumerate(slides_data):
                self._create_slide(prs, slide_info, idx)
                result["slide_count"] += 1

            file_name = f"{task_id}_{topic.replace(' ', '_')}.pptx"
            file_path = self.output_dir / file_name
            prs.save(str(file_path))

            result["success"] = True
            result["file_path"] = str(file_path)
            result["file_url"] = f"/static/pptx/{file_name}"
            result["download_name"] = f"{topic}_演示文稿.pptx"

            logger.info(
                "✅ PPTX 生成成功: %s (%d页)",
                file_name, result["slide_count"]
            )

        except ImportError as e:
            result["error"] = f"python-pptx 未安装: {e}。请运行: pip install python-pptx"
            logger.error("❌ python-pptx 未安装: %s", e)
        except Exception as e:
            result["error"] = f"PPTX 生成异常: {str(e)}"
            logger.error("❌ PPTX 生成异常: %s", e, exc_info=True)

        return result

    # ── Markdown 解析 ──

    def _parse_markdown(self, content: str, title: str) -> List[Dict]:
        slides = []
        lines = content.split("\n")
        current_slide = None
        in_code_block = False
        code_lines = []

        # 封面页
        slides.append({"layout": "title", "title": title})

        # 用于清理 bullet 中 LLM 残留的元信息前缀
        def _clean_bullet(text: str) -> str:
            text = text.strip()
            # 去掉 markdown 加粗
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
            # 去掉"要点 N:"、"核心要点 N:" 等前缀
            text = re.sub(r'^(?:核心)?要点\s*\d*[：:]\s*', '', text)
            text = re.sub(r'^(?:关键)?图示[：:].*', '', text)
            text = re.sub(r'^演讲备注[：:].*', '', text)
            text = re.sub(r'^视觉提示[：:].*', '', text)
            text = re.sub(r'^副标题[：:]\s*', '', text)
            return text.strip()

        for line in lines:
            line_stripped = line.strip()

            # 代码块处理
            if line_stripped.startswith("```"):
                if in_code_block and current_slide is not None:
                    current_slide["code_block"] = "\n".join(code_lines)
                    code_lines = []
                in_code_block = not in_code_block
                continue

            if in_code_block:
                code_lines.append(line)
                continue

            # 跳过空行和分隔线
            if not line_stripped or line_stripped == "---":
                continue

            # 一级标题 → 章节页
            if line_stripped.startswith("# ") and not line_stripped.startswith("##"):
                if current_slide:
                    slides.append(current_slide)
                h1_text = line_stripped[2:].strip()
                current_slide = {"layout": "section", "title": h1_text}
                continue

            # 二级标题 → 新幻灯片
            if line_stripped.startswith("## "):
                if current_slide:
                    slides.append(current_slide)
                h2_text = line_stripped[3:].strip()
                # 检测是否为代码页
                is_code_page = any(kw in h2_text for kw in ["代码", "Code", "code", "示例代码"])
                current_slide = {
                    "layout": "code" if is_code_page else "content",
                    "title": h2_text,
                    "bullets": [],
                }
                continue

            # 列表项
            if line_stripped.startswith(("- ", "* ", "+ ")):
                bullet_text = re.sub(r"^[-*+]\s+", "", line_stripped)
                bullet_text = _clean_bullet(bullet_text)
                if not bullet_text:
                    continue
                if current_slide is None:
                    current_slide = {"layout": "content", "title": "", "bullets": []}
                current_slide.setdefault("bullets", []).append(bullet_text)
                continue

            # 编号列表 (1. xxx)
            if re.match(r'^\d+[\.\)]\s', line_stripped):
                bullet_text = re.sub(r'^\d+[\.\)]\s+', '', line_stripped)
                bullet_text = _clean_bullet(bullet_text)
                if not bullet_text:
                    continue
                if current_slide is None:
                    current_slide = {"layout": "content", "title": "", "bullets": []}
                current_slide.setdefault("bullets", []).append(bullet_text)
                continue

            # 引用块
            if line_stripped.startswith("> "):
                quote_text = line_stripped[2:].strip()
                if current_slide is not None:
                    current_slide.setdefault("quotes", []).append(quote_text)
                continue

            # 其他文本 → 追加到 extra_text
            if current_slide is not None and line_stripped:
                current_slide.setdefault("extra_text", "")
                current_slide["extra_text"] += line_stripped + "\n"

        if current_slide:
            slides.append(current_slide)

        # 兜底：如果解析不出内容，生成降级页
        if len(slides) <= 1:
            text_content = re.sub(r'#+\s*', '', content)
            fallback_bullets = [
                _clean_bullet(p.strip())
                for p in text_content.split('\n')
                if p.strip()
            ][:6]
            slides.append({
                "layout": "content",
                "title": title,
                "bullets": fallback_bullets or ["请参考详细内容"],
            })

        return slides[:15]

    # ── 形状工具方法 ──

    @staticmethod
    def _add_rect(slide, left, top, width, height, color_hex):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        shape.line.fill.background()
        return shape

    @staticmethod
    def _add_oval(slide, left, top, size, color_hex):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        shape.line.fill.background()
        return shape

    @staticmethod
    def _add_rounded_rect(slide, left, top, width, height, color_hex):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        shape.line.fill.background()

        # 圆角调整
        try:
            shape.adjustments[0] = 0.15
        except (IndexError, AttributeError):
            pass
        return shape

    @staticmethod
    def _add_text(slide, left, top, width, height, text,
                  font_name="Microsoft YaHei", font_size=18,
                  font_bold=False, color_hex="1D3557",
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.TOP):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = font_bold
        p.font.name = font_name
        p.font.color.rgb = RGBColor.from_string(color_hex)
        p.alignment = align
        try:
            tf.paragraphs[0].space_before = Pt(0)
            tf.paragraphs[0].space_after = Pt(0)
        except Exception:
            pass
        return txBox

    # ════════════════════════════════════════
    #   页面构建器（科大讯飞风格）
    # ════════════════════════════════════════

    def _create_slide(self, prs, info: Dict, index: int):
        layout_type = info.get("layout", "content")
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        builders = {
            "title": self._build_title_slide,
            "section": self._build_section_slide,
            "content": self._build_content_slide,
            "code": self._build_code_slide,
        }

        builder = builders.get(layout_type, self._build_content_slide)
        builder(slide, info, index)
        return slide

    def _build_title_slide(self, slide, info, index=0):
        """
        封面页 — 科大讯飞风格：
        - 左上角品牌 Logo 区域
        - 居中大标题（深蓝/红色）
        - 底部装饰波浪线/渐变条
        - 右下角 AI 标识
        """
        t = THEME
        W = Inches(13.333)
        H = Inches(7.5)

        # 白色底（默认）
        # 左侧品牌色竖条
        self._add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, t["brand_red"])

        # 品牌标识区域（左上）
        self._add_text(
            slide, Inches(0.6), Inches(0.5), Inches(3), Inches(0.6),
            "StructQuest AI", "Microsoft YaHei", 16, True,
            t["brand_red"], PP_ALIGN.LEFT
        )

        # 副标签
        self._add_text(
            slide, Inches(0.6), Inches(1.0), Inches(4), Inches(0.4),
            "智能学习平台 · 演示文稿", "Microsoft YaHei", 12, False,
            t["text_light"], PP_ALIGN.LEFT
        )

        # 主标题（居中偏下，醒目位置）
        title_text = info.get("title", "演示文稿")
        self._add_text(
            slide, Inches(1), Inches(2.4), Inches(11.333), Inches(1.8),
            title_text, t["title_font"], t["title_size"], True,
            t["brand_blue"], PP_ALIGN.CENTER
        )

        # 装饰线（标题下方）
        self._add_rect(
            slide, Inches(4.5), Inches(4.4), Inches(4.333), Inches(0.06),
            t["accent_teal"]
        )

        # 底部信息栏（浅灰背景）
        self._add_rect(slide, Inches(0), Inches(6.6), W, Inches(0.9), t["bg_light"])

        # 底部文字
        self._add_text(
            slide, Inches(0.8), Inches(6.75), Inches(6), Inches(0.5),
            "AI 驱动 · 智能生成 · 专业设计", "Microsoft YaHei",
            14, False, t["text_gray"], PP_ALIGN.LEFT
        )

        # 右下角 AI 标识
        self._add_text(
            slide, Inches(10), Inches(6.75), Inches(2.8), Inches(0.5),
            "AI | 用人工智能建设美好世界", "Microsoft YaHei",
            11, False, t["text_light"], PP_ALIGN.RIGHT
        )

    def _build_section_slide(self, slide, info, index=0):
        """
        章节分隔页 — 全屏品牌色背景 + 居中白色大字
        类似科大讯飞的章节过渡效果
        """
        t = THEME
        W = Inches(13.333)
        H = Inches(7.5)

        # 根据索引轮换不同的强调色
        colors = [t["brand_red"], t["accent_teal"], t["accent_blue"],
                   t["accent_purple"], t["accent_orange"]]
        bg_color = colors[index % len(colors)]

        # 全屏色块背景
        self._add_rect(slide, Inches(0), Inches(0), W, H, bg_color)

        # 装饰性半透明圆形（右上角）
        self._add_oval(slide, Inches(9.5), Inches(-1.5), Inches(4), "FFFFFF")
        # 降低透明度通过覆盖一层
        oval_shape = slide.shapes[-1]
        try:
            oval_shape.fill.fore_color.rgb.__class__  # just check
        except Exception:
            pass

        # 章节标题
        section_title = info.get("title", "")
        self._add_text(
            slide, Inches(1), Inches(2.8), Inches(11.333), Inches(2),
            section_title, t["title_font"], 44, True,
            t["white"], PP_ALIGN.CENTER
        )

        # 装饰短线
        self._add_rect(
            slide, Inches(5.667), Inches(5.0), Inches(2), Inches(0.05),
            t["white"]
        )

        # 页码提示
        self._add_text(
            slide, Inches(1), Inches(5.4), Inches(11.333), Inches(0.5),
            "— 章节 —", "Microsoft YaHei", 14, False,
            "FFCCCC", PP_ALIGN.CENTER
        )

    def _build_content_slide(self, slide, info, index=0):
        """
        内容页 — 全视觉卡片布局

        所有要点都用彩色卡片展示：
        - 1~4 个要点: 水平一排大卡片（带圆形图标）
        - 5~6 个要点: 2×3 网格中等卡片
        - 7+ 个要点: 3 列网格紧凑卡片
        """
        t = THEME
        bullets = info.get("bullets", [])
        page_title = info.get("title", "")
        extra = info.get("extra_text", "").strip()
        quotes = info.get("quotes", [])

        # 顶部标题栏
        title_color = t["circle_colors"][index % len(t["circle_colors"])][0]
        self._add_rect(
            slide, Inches(0), Inches(0), Inches(13.333), Inches(1.15),
            title_color
        )
        self._add_text(
            slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.8),
            page_title, t["title_font"], t["heading1_size"],
            True, t["white"], PP_ALIGN.LEFT
        )

        clean_bullets = [b for b in bullets if b and len(b) > 0]

        if clean_bullets:
            n = len(clean_bullets)
            if n <= 4:
                self._render_circle_cards(slide, clean_bullets)
            else:
                self._render_card_grid(slide, clean_bullets)
        elif quotes:
            self._render_quotes(slide, quotes)
        elif extra:
            self._add_text(
                slide, Inches(0.8), Inches(1.45), Inches(11.7), Inches(5.2),
                extra[:600], t["body_font"], t["body_size"],
                False, t["text_gray"], PP_ALIGN.LEFT
            )

        # 页脚
        self._add_rect(slide, Inches(0), Inches(7.15), Inches(13.333),
                       Inches(0.35), t["bg_light"])
        self._add_text(
            slide, Inches(11.5), Inches(7.18), Inches(1.6), Inches(0.28),
            "StructQuest AI", "Microsoft YaHei", 9, False,
            t["text_light"], PP_ALIGN.RIGHT
        )

    def _render_circle_cards(self, slide, bullets):
        """
        渲染彩色圆圈卡片布局（科大讯飞风格核心）

        每个要点显示为：
          ┌────────────────────┐
          │     ⭕ 彩色圆圈     │
          │     中文标题       │
          │   English Subtitle │
          │   ─ ─ ─ ─ ─ ─ ─  │
          │   详细描述文字...  │
          └────────────────────┘
        """
        t = THEME
        n = len(bullets)
        card_w = Inches(2.95)
        card_h = Inches(5.0)
        gap = Inches(0.35)
        start_x = Inches(0.65)
        start_y = Inches(1.45)

        # 计算起始 X 使所有卡片居中
        total_w = n * card_w + (n - 1) * gap
        if total_w < Inches(12):
            start_x = (Inches(13.333) - total_w) / 2

        for i, bullet in enumerate(bullets):
            x = start_x + i * (card_w + gap)
            color_pair = t["circle_colors"][i % len(t["circle_colors"])]
            main_color = color_pair[0]

            # 卡片背景（白色带阴影效果用边框模拟）
            card = self._add_rounded_rect(
                slide, x, start_y, card_w, card_h, t["white"]
            )
            try:
                card.line.color.rgb = __import__(
                    "pptx.dml.color", fromlist=["RGBColor"]
                ).RGBColor.from_string(t["card_shadow"])
                card.line.width = Pt(1)
            except Exception:
                pass

            # 解析 bullet 文本
            lines = self._split_bullet_text(bullet)
            card_title = lines[0] if lines else bullet
            subtitle = ""
            desc_lines = []

            if len(lines) >= 2:
                # 第二行可能是副标题或描述
                if len(lines[1]) < 25 and not lines[1].startswith(("例如", "包括", "如")):
                    subtitle = lines[1]
                    desc_lines = lines[2:]
                else:
                    desc_lines = lines[1:]

            desc_text = "\n".join(desc_lines[:4]) if desc_lines else ""

            # ── 圆形色块 ──
            circle_size = Inches(1.4)
            circle_x = x + (card_w - circle_size) / 2
            circle_y = start_y + Inches(0.3)
            self._add_oval(slide, circle_x, circle_y, circle_size, main_color)

            # 圆圈内文字（首字或图标占位）
            icon_char = card_title[0] if card_title else "?"
            self._add_text(
                slide, circle_x, circle_y + Inches(0.42),
                circle_size, Inches(0.6),
                icon_char.upper(), "Microsoft YaHei",
                28, True, t["white"], PP_ALIGN.CENTER
            )

            # ── 卡片标题 ──
            title_y = start_y + Inches(1.85)
            self._add_text(
                slide, x + Inches(0.12), title_y,
                card_w - Inches(0.24), Inches(0.55),
                card_title, t["body_font"], 17,
                True, t["text_dark"], PP_ALIGN.CENTER
            )

            # ── 英文副标题（如有） ──
            if subtitle:
                sub_y = title_y + Inches(0.55)
                self._add_text(
                    slide, x + Inches(0.12), sub_y,
                    card_w - Inches(0.24), Inches(0.38),
                    subtitle, "Arial", 12, False,
                    main_color, PP_ALIGN.CENTER
                )

            # ── 分隔虚线 ──
            dash_y = start_y + Inches(2.95) if subtitle else start_y + Inches(2.55)
            self._add_rect(
                slide, x + Inches(0.5), dash_y,
                card_w - Inches(1.0), Inches(0.02),
                t["card_shadow"]
            )

            # ── 描述文字 ──
            if desc_text:
                desc_y = dash_y + Inches(0.2)
                self._add_text(
                    slide, x + Inches(0.15), desc_y,
                    card_w - Inches(0.3), Inches(1.9),
                    desc_text, t["body_font"], t["body_size"],
                    False, t["text_gray"], PP_ALIGN.CENTER
                )

    def _render_card_grid(self, slide, bullets):
        """
        网格卡片布局 — 适合 5+ 个要点

        自动计算行列：2列或3列网格，每格一张彩色卡片
        """
        t = THEME
        n = len(bullets)
        if n <= 6:
            cols, rows = 2, (n + 1) // 2
        else:
            cols, rows = 3, (n + 2) // 3

        card_w = Inches(12.0 / cols - 0.25)
        card_h = Inches(5.3 / rows - 0.15)
        start_x = Inches(0.65)
        start_y = Inches(1.45)
        gap_x = Inches(0.25)
        gap_y = Inches(0.15)

        for i, bullet in enumerate(bullets):
            row = i // cols
            col = i % cols
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)

            color_pair = t["circle_colors"][i % len(t["circle_colors"])]
            bg_color = color_pair[0]

            # 卡片背景
            self._add_rounded_rect(slide, x, y, card_w, card_h, t["white"])
            try:
                from pptx.util import Pt as _Pt
                slide.shapes[-1].line.color.rgb = RGBColor.from_string(t["card_shadow"])
                slide.shapes[-1].line.width = _Pt(1)
            except Exception:
                pass

            # 左侧色条
            self._add_rect(
                slide, x, y, Inches(0.08), card_h, bg_color
            )

            # 文字内容
            lines = self._split_bullet_text(bullet, max_len=20)
            main_text = lines[0] if lines else bullet
            sub_text = ""
            if len(lines) >= 2:
                sub_text = lines[1][:60]

            # 标题（大号、加粗、彩色）
            self._add_text(
                slide, x + Inches(0.2), y + Inches(0.15),
                card_w - Inches(0.35), Inches(0.5),
                main_text[:30], t["body_font"], 17,
                True, bg_color, PP_ALIGN.LEFT
            )

            # 副文本（小号灰色）
            if sub_text:
                self._add_text(
                    slide, x + Inches(0.2), y + Inches(0.65),
                    card_w - Inches(0.35), card_h - Inches(0.8),
                    sub_text, t["body_font"], 13,
                    False, t["text_gray"], PP_ALIGN.LEFT
                )

    def _render_quotes(self, slide, quotes):
        """对比/引用页 — 左右分栏展示"""
        t = THEME
        n = len(quotes)
        if n == 2:
            # 左右对比
            col_w = Inches(5.8)
            for i, quote in enumerate(quotes[:2]):
                x = Inches(0.8) + i * Inches(6.2)
                self._add_rounded_rect(
                    slide, x, Inches(1.6), col_w, Inches(4.5), t["white"]
                )
                # 色条顶部
                color = t["circle_colors"][i % len(t["circle_colors"])][0]
                self._add_rect(
                    slide, x, Inches(1.6), col_w, Inches(0.08), color
                )
                self._add_text(
                    slide, x + Inches(0.3), Inches(1.9),
                    col_w - Inches(0.6), Inches(3.8),
                    quote[:300], t["body_font"], 14,
                    False, t["text_dark"], PP_ALIGN.LEFT
                )
        else:
            # 单列引用
            for i, quote in enumerate(quotes[:5]):
                y = Inches(1.6) + i * Inches(1.1)
                color = t["circle_colors"][i % len(t["circle_colors"])][0]
                self._add_rect(
                    slide, Inches(0.8), y, Inches(0.08), Inches(0.8), color
                )
                self._add_text(
                    slide, Inches(1.1), y, Inches(11.5), Inches(0.9),
                    quote[:200], t["body_font"], 15,
                    False, t["text_dark"], PP_ALIGN.LEFT
                )

    def _build_code_slide(self, slide, info, index=0):
        """代码展示页 — 深色背景 + 代码高亮"""
        t = THEME
        page_title = info.get("title", "代码示例")
        code_block = info.get("code_block", "")
        bullets = info.get("bullets", [])

        W = Inches(13.333)

        # 深色标题栏
        self._add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), t["brand_blue"])
        self._add_text(
            slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.8),
            page_title, t["title_font"], t["heading1_size"],
            True, t["white"], PP_ALIGN.LEFT
        )

        if code_block:
            # 代码区域深色背景
            code_bg_color = "1E1E2E"
            self._add_rect(
                slide, Inches(0.6), Inches(1.5),
                Inches(8.5), Inches(5.3),
                code_bg_color
            )
            # 代码文字（亮色）
            code_lines = code_block.strip().split("\n")[:20]
            code_text = "\n".join(code_lines)
            self._add_text(
                slide, Inches(0.85), Inches(1.7),
                Inches(8.1), Inches(5.0),
                code_text, t["code_font"], 13,
                False, "E0E0E0", PP_ALIGN.LEFT
            )

        # 右侧说明（如果有 bullets）
        if bullets:
            for i, bullet in enumerate(bullets[:4]):
                y = Inches(1.5) + Inches(i * 1.2)
                color = t["circle_colors"][i % len(t["circle_colors"])][0]
                self._add_rect(
                    slide, Inches(9.6), y + Inches(0.05),
                    Inches(0.08), Inches(0.7), color
                )
                self._add_text(
                    slide, Inches(9.9), y, Inches(3.1), Inches(0.85),
                    bullet[:80], t["body_font"], 14,
                    False, t["text_dark"], PP_ALIGN.LEFT
                )

        # 页脚
        self._add_rect(slide, Inches(0), Inches(7.15), W, Inches(0.35), t["bg_light"])
        self._add_text(
            slide, Inches(11.5), Inches(7.18), Inches(1.6), Inches(0.28),
            "StructQuest AI", "Microsoft YaHei", 9, False,
            t["text_light"], PP_ALIGN.RIGHT
        )

    @staticmethod
    def _split_bullet_text(text: str, max_len: int = 50) -> List[str]:
        """将一条 bullet 文本拆分为标题+描述行"""
        # 尝试按常见分隔符拆分
        for sep in ["：", ": ", "——", "--", "。"]:
            if sep in text:
                parts = text.split(sep, 1)
                if len(parts[0]) < max_len:
                    return [parts[0].strip(), sep.strip() + parts[1].strip()]

        # 没有明确分隔符，按长度切分
        if len(text) <= max_len:
            return [text]

        # 找第一个句号或换行
        for cut_ch in ["。", "；", ";", "."]:
            idx = text.find(cut_ch)
            if 5 < idx < max_len:
                return [text[:idx + cut_ch].strip(), text[idx + len(cut_ch):].strip()]

        return [text[:max_len], text[max_len:]]


# ── 全局单例 ──

_instance = None


def get_pptx_generator() -> PPTXGenerator:
    global _instance
    if _instance is None:
        _instance = PPTXGenerator()
    return _instance
